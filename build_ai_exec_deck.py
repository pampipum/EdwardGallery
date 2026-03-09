from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()

# 16:9 widescreen for modern displays
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_title_slide(title, subtitle, notes):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    notes_frame = slide.notes_slide.notes_text_frame
    notes_frame.text = notes


def add_bullet_slide(title, bullets, notes):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()

    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(24)

    notes_frame = slide.notes_slide.notes_text_frame
    notes_frame.text = notes


add_title_slide(
    "AI in 2026: Practical Adoption for Swiss Mid-Size Enterprises",
    "Executive Briefing | March 2026",
    "Set the tone: execution over hype. Emphasize business outcomes, risk discipline, and realistic timelines."
)

slides = [
    (
        "Executive Thesis",
        [
            "AI value now depends on operating model, not model novelty",
            "Winners combine fast pilots with governance-by-design",
            "Focus on workflow redesign, not chatbot deployment",
        ],
        "Use German snippet if useful: 'Entscheidungsreife statt Hype.'"
    ),
    (
        "Why Now (2025 -> 2026)",
        [
            "Model capability is sufficient for enterprise-grade use cases",
            "Cost/performance improved; choice increased across vendors",
            "Competitors are moving from pilots to operating models",
        ],
        "Stress urgency without panic: waiting 12 months increases catch-up cost."
    ),
    (
        "AI Economics (CFO View)",
        [
            "Measure cost per workflow outcome, not only per-token cost",
            "Track quality-adjusted productivity and rework reduction",
            "Use stage-gates to stop weak pilots quickly",
        ],
        "Introduce KPI set: cycle time, error rate, unit cost, compliance incidents, adoption rate."
    ),
    (
        "Frontier Landscape Snapshot",
        [
            "Closed ecosystems: strongest managed tooling and enterprise packaging",
            "Open-weight ecosystem: flexibility, lower lock-in, controllable deployment",
            "Multi-model architecture is now a practical default",
        ],
        "Position this as portfolio design, not brand preference."
    ),
    (
        "Open vs Closed: Decision Lens",
        [
            "Performance on your tasks and languages (EN/DE/FR)",
            "Compliance posture: data handling, traceability, controls",
            "Portability and 24-month total cost of ownership",
        ],
        "Recommend benchmark harness on internal tasks before committing at scale."
    ),
    (
        "Chinese Model Ecosystem",
        [
            "Rapid open-weight progress (e.g., DeepSeek, Qwen families)",
            "Stronger buyer leverage on price/performance globally",
            "Higher due diligence needs: provenance, security, jurisdiction",
        ],
        "Avoid geopolitics-heavy framing; keep it procurement and risk focused."
    ),
    (
        "The Agents Wave",
        [
            "Copilot -> Agent -> Multi-agent workflow orchestration",
            "Agents add tool use, memory, planning, and handoffs",
            "Value rises with guardrails, observability, and fallback paths",
        ],
        "Clarify that agents are not autonomous by default; they are controllable systems."
    ),
    (
        "OpenClaw Practical Pattern",
        [
            "Agent loop with explicit tools and state",
            "Human approvals for external/high-risk actions",
            "Auditable actions and repeatable runbooks",
        ],
        "Use as a practical architecture example: controlled autonomy in enterprise settings."
    ),
    (
        "Enterprise Pattern 1: Copilots",
        [
            "Fastest time-to-value for individual productivity",
            "Low integration burden; easy to launch",
            "Risk: limited process impact if not connected to workflows",
        ],
        "Use copilots as on-ramp, not end-state."
    ),
    (
        "Enterprise Pattern 2: RAG",
        [
            "Grounded answers from internal documents and policies",
            "Effectiveness depends on data quality and permissions",
            "Best for knowledge-heavy support and decision preparation",
        ],
        "Call out need for citation, freshness checks, and access controls."
    ),
    (
        "Enterprise Pattern 3: Agentic Workflows",
        [
            "Automates multi-step work across systems",
            "Highest upside for operations and back-office processes",
            "Requires strict testing, rollback, and monitoring",
        ],
        "Encourage narrow scope first, then scale once quality thresholds are met."
    ),
    (
        "AI Operating Model",
        [
            "Business product owner accountable for outcomes",
            "AI platform team provides shared capabilities",
            "Risk/compliance embedded from design to production",
        ],
        "Recommend monthly governance forum and quarterly value review."
    ),
    (
        "EU and Swiss Governance",
        [
            "EU AI Act obligations phased across 2025-2027",
            "Swiss FADP already applies to AI data processing",
            "Implement controls now; do not wait for perfect legal clarity",
        ],
        "Keep legal details in appendix; emphasize practical controls in main story."
    ),
    (
        "Human Oversight by Design",
        [
            "Human-in-the-loop for high-impact decisions",
            "Human-on-the-loop for scaled low-risk automation",
            "Defined escalation thresholds and override authority",
        ],
        "Define where humans must approve before action (finance, legal, customer commitments)."
    ),
    (
        "AI Quality Management",
        [
            "Pre-production evaluations on real tasks",
            "Regression tests after model or prompt changes",
            "Error budgets and hallucination tolerance per workflow",
        ],
        "Treat model updates like software releases with QA gates."
    ),
    (
        "Security and Trust",
        [
            "Defend against prompt injection and data exfiltration",
            "Protect against deepfake-enabled fraud in operations",
            "Maintain audit logs with privacy-preserving redaction",
        ],
        "Security leader should co-own deployment standards with AI platform lead."
    ),
    (
        "Data Readiness: Hidden Bottleneck",
        [
            "Data classification and access controls",
            "Metadata, lineage, and document lifecycle hygiene",
            "Policy-aligned retention and deletion",
        ],
        "State clearly: poor data readiness is the #1 reason RAG pilots underperform."
    ),
    (
        "Build vs Buy Framework",
        [
            "Buy for commodity capabilities and speed",
            "Build/customize for differentiated workflows and IP",
            "Keep model layer portable to reduce lock-in",
        ],
        "Use 6 gates: differentiation, risk, speed, integration, talent, TCO."
    ),
    (
        "30/60/90 Day Adoption Roadmap",
        [
            "30: governance baseline, benchmark harness, use-case selection",
            "60: two pilots with production-like security and metrics",
            "90: scale one winner, stop one loser, formalize operating model",
        ],
        "German snippet option: 'Nächster Schritt: Entscheidung in 2 Wochen.'"
    ),
    (
        "Use-Case Blueprint 1: Customer Service",
        [
            "AI triage + response drafting + knowledge retrieval",
            "Human escalation for sensitive customer cases",
            "KPIs: first response time, deflection, CSAT",
        ],
        "Ideal early use case due clear volume and measurable impact."
    ),
    (
        "Use-Case Blueprint 2: Sales and RFP",
        [
            "RAG-driven proposal drafting and evidence retrieval",
            "Compliance/legal checks before external submission",
            "KPIs: cycle time, bid quality, win-rate contribution",
        ],
        "Good fit for multilingual Swiss market dynamics."
    ),
    (
        "Use-Case Blueprint 3: Procure-to-Pay",
        [
            "Invoice/contract extraction and policy checks",
            "Exception routing with segregation of duties",
            "KPIs: processing time, error rate, audit exceptions",
        ],
        "Strong case for controlled agentic automation in finance ops."
    ),
    (
        "Executive Decisions Required",
        [
            "Approve top 3 use cases and accountable owners",
            "Approve governance baseline and risk thresholds",
            "Approve architecture principle: multi-model, portable, auditable",
        ],
        "Close with concrete asks and timeline."
    ),
]

for title, bullets, notes in slides:
    add_bullet_slide(title, bullets, notes)

out_path = "/root/.openclaw/workspace/AI_Executive_Briefing_Zurich_2026.pptx"
prs.save(out_path)
print(out_path)
