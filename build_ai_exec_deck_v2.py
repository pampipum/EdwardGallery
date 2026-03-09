from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Palette: Midnight Executive
C_NAVY = RGBColor(0x1E, 0x27, 0x61)
C_ICE = RGBColor(0xCA, 0xDC, 0xFC)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT = RGBColor(0xF6, 0xF8, 0xFC)
C_TEXT = RGBColor(0x1F, 0x29, 0x3A)
C_ACCENT = RGBColor(0x2E, 0xB8, 0xA6)
C_MUTED = RGBColor(0x63, 0x6E, 0x85)

FONT_TITLE = "Calibri"
FONT_BODY = "Calibri"


def set_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def full_bg(slide, color):
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)


def title_dark(title, subtitle, notes):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    full_bg(slide, C_NAVY)

    # motif panel
    panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(8.6), Inches(0.7), Inches(4.1), Inches(6.1))
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor(0x2A, 0x35, 0x7A)
    panel.line.fill.background()

    ring = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.DONUT, Inches(9.6), Inches(2.0), Inches(2.2), Inches(2.2))
    ring.fill.solid()
    ring.fill.fore_color.rgb = C_ACCENT
    ring.line.fill.background()

    tx = slide.shapes.add_textbox(Inches(0.9), Inches(1.4), Inches(7.4), Inches(2.6)).text_frame
    tx.clear()
    p = tx.paragraphs[0]
    p.text = title
    p.font.name = FONT_TITLE
    p.font.bold = True
    p.font.size = Pt(44)
    p.font.color.rgb = C_WHITE

    sp = tx.add_paragraph()
    sp.text = subtitle
    sp.font.name = FONT_BODY
    sp.font.size = Pt(20)
    sp.font.color.rgb = C_ICE
    sp.space_before = Pt(14)

    foot = slide.shapes.add_textbox(Inches(0.9), Inches(6.5), Inches(8.5), Inches(0.5)).text_frame
    foot.paragraphs[0].text = "Executive briefing • Zurich • March 2026"
    foot.paragraphs[0].font.name = FONT_BODY
    foot.paragraphs[0].font.size = Pt(12)
    foot.paragraphs[0].font.color.rgb = C_ICE

    set_notes(slide, notes)


def title_light(title, bullets, notes):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    full_bg(slide, C_LIGHT)

    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.2))
    band.fill.solid()
    band.fill.fore_color.rgb = C_NAVY
    band.line.fill.background()

    tbox = slide.shapes.add_textbox(Inches(0.8), Inches(0.22), Inches(11.7), Inches(0.8)).text_frame
    tbox.paragraphs[0].text = title
    tbox.paragraphs[0].font.name = FONT_TITLE
    tbox.paragraphs[0].font.bold = True
    tbox.paragraphs[0].font.size = Pt(32)
    tbox.paragraphs[0].font.color.rgb = C_WHITE

    y = 1.8
    for i, b in enumerate(bullets):
        card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(y), Inches(11.5), Inches(1.35))
        card.fill.solid()
        card.fill.fore_color.rgb = C_WHITE
        card.line.color.rgb = C_ICE
        tf = card.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = b
        p.font.name = FONT_BODY
        p.font.size = Pt(21)
        p.font.color.rgb = C_TEXT
        y += 1.65

    dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(12.2), Inches(6.6), Inches(0.35), Inches(0.35))
    dot.fill.solid()
    dot.fill.fore_color.rgb = C_ACCENT
    dot.line.fill.background()

    set_notes(slide, notes)


def two_col(title, left_title, left_points, right_title, right_points, notes):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    full_bg(slide, C_LIGHT)

    tt = slide.shapes.add_textbox(Inches(0.8), Inches(0.35), Inches(11.8), Inches(0.8)).text_frame
    p = tt.paragraphs[0]
    p.text = title
    p.font.name = FONT_TITLE
    p.font.bold = True
    p.font.size = Pt(34)
    p.font.color.rgb = C_NAVY

    left = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.4), Inches(5.85), Inches(5.6))
    left.fill.solid()
    left.fill.fore_color.rgb = C_WHITE
    left.line.color.rgb = C_ICE
    right = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(6.7), Inches(1.4), Inches(5.85), Inches(5.6))
    right.fill.solid()
    right.fill.fore_color.rgb = C_WHITE
    right.line.color.rgb = C_ICE

    ltf = left.text_frame
    ltf.clear()
    lp = ltf.paragraphs[0]
    lp.text = left_title
    lp.font.bold = True
    lp.font.size = Pt(23)
    lp.font.color.rgb = C_NAVY
    for item in left_points:
        p = ltf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(17)
        p.font.color.rgb = C_TEXT

    rtf = right.text_frame
    rtf.clear()
    rp = rtf.paragraphs[0]
    rp.text = right_title
    rp.font.bold = True
    rp.font.size = Pt(23)
    rp.font.color.rgb = C_NAVY
    for item in right_points:
        p = rtf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(17)
        p.font.color.rgb = C_TEXT

    set_notes(slide, notes)


def cards3(title, cards, notes):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    full_bg(slide, C_LIGHT)

    t = slide.shapes.add_textbox(Inches(0.8), Inches(0.35), Inches(11.5), Inches(0.8)).text_frame
    p = t.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(34)
    p.font.color.rgb = C_NAVY

    x = 0.8
    for heading, lines in cards:
        card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.5), Inches(4.05), Inches(5.4))
        card.fill.solid()
        card.fill.fore_color.rgb = C_WHITE
        card.line.color.rgb = C_ICE

        icon = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x + 0.25), Inches(1.75), Inches(0.45), Inches(0.45))
        icon.fill.solid()
        icon.fill.fore_color.rgb = C_ACCENT
        icon.line.fill.background()

        tf = card.text_frame
        tf.clear()
        hp = tf.paragraphs[0]
        hp.text = f"   {heading}"
        hp.font.bold = True
        hp.font.size = Pt(22)
        hp.font.color.rgb = C_NAVY
        for ln in lines:
            p = tf.add_paragraph()
            p.text = f"• {ln}"
            p.font.size = Pt(16)
            p.font.color.rgb = C_TEXT
        x += 4.25

    set_notes(slide, notes)


def timeline(title, steps, notes):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    full_bg(slide, C_LIGHT)

    t = slide.shapes.add_textbox(Inches(0.8), Inches(0.35), Inches(11.6), Inches(0.8)).text_frame
    t.paragraphs[0].text = title
    t.paragraphs[0].font.bold = True
    t.paragraphs[0].font.size = Pt(34)
    t.paragraphs[0].font.color.rgb = C_NAVY

    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(1.3), Inches(3.55), Inches(10.8), Inches(0.08))
    line.fill.solid()
    line.fill.fore_color.rgb = C_NAVY
    line.line.fill.background()

    x = 1.2
    for i, (phase, detail) in enumerate(steps, start=1):
        circ = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(3.2), Inches(0.75), Inches(0.75))
        circ.fill.solid()
        circ.fill.fore_color.rgb = C_ACCENT
        circ.line.fill.background()

        n = slide.shapes.add_textbox(Inches(x+0.23), Inches(3.37), Inches(0.3), Inches(0.2)).text_frame
        n.paragraphs[0].text = str(i)
        n.paragraphs[0].font.bold = True
        n.paragraphs[0].font.size = Pt(16)
        n.paragraphs[0].font.color.rgb = C_WHITE

        top = slide.shapes.add_textbox(Inches(x-0.35), Inches(2.0), Inches(1.6), Inches(1.0)).text_frame
        top.paragraphs[0].text = phase
        top.paragraphs[0].font.bold = True
        top.paragraphs[0].font.size = Pt(18)
        top.paragraphs[0].font.color.rgb = C_NAVY
        top.paragraphs[0].alignment = PP_ALIGN.CENTER

        bot = slide.shapes.add_textbox(Inches(x-0.8), Inches(4.15), Inches(2.5), Inches(1.2)).text_frame
        bot.paragraphs[0].text = detail
        bot.paragraphs[0].font.size = Pt(14)
        bot.paragraphs[0].font.color.rgb = C_TEXT
        bot.paragraphs[0].alignment = PP_ALIGN.CENTER

        x += 3.4

    set_notes(slide, notes)


title_dark(
    "AI in 2026: Practical Adoption for Swiss Mid-Size Enterprises",
    "Executive presentation for a Zurich business audience",
    "Set tone: practical, no-hype, decision-useful."
)

title_light(
    "Executive Thesis",
    [
        "AI advantage now comes from execution discipline, not model novelty",
        "Use portfolio logic: pilots, stage-gates, scale winners quickly",
        "Design governance and security into architecture from day one",
    ],
    "Optional German: Entscheidungsreife statt Hype."
)

two_col(
    "What Changed (2025-2026)",
    "Technology Shift",
    ["Tool-capable models improved", "Open-weight quality increased", "Model costs continue to decline"],
    "Business Shift",
    ["AI adoption widened", "Agent pilots became mainstream", "Boards ask for measurable ROI"],
    "Explain that strategy must update from experimentation to operating model."
)

cards3(
    "Frontier Model Landscape",
    [
        ("Closed Ecosystems", ["Managed reliability", "Integrated tooling", "Enterprise packaging"]),
        ("Open-Weight", ["Portability & control", "Cost leverage", "Faster customization"]),
        ("Practical Reality", ["No single winner", "Use workload routing", "Benchmark on your data"]),
    ],
    "Emphasize multi-model architecture as default in 2026."
)

two_col(
    "Open vs Closed: Decision Framework",
    "Buy Closed When",
    ["Speed matters most", "Managed controls are sufficient", "Internal AI team is small"],
    "Prefer Open/Hybrid When",
    ["Data control is critical", "Need deep customization", "Vendor lock-in is unacceptable"],
    "This is a business decision, not an ideological one."
)

cards3(
    "Chinese Model Ecosystem: Why It Matters",
    [
        ("Capability", ["Rapid progress from Qwen/DeepSeek families", "Global benchmark gaps narrowing", "Stronger competitive pressure"]),
        ("Procurement", ["Better buyer leverage", "More model choices", "Need stricter due diligence"]),
        ("Risk Lens", ["Provenance checks", "Security controls", "Jurisdiction-aware architecture"]),
    ],
    "Keep framing practical: impact on pricing, risk, and architecture decisions."
)

cards3(
    "The Agent Wave",
    [
        ("From Copilot", ["Suggests content", "Supports human tasks", "Limited automation"]),
        ("To Agent", ["Uses tools", "Handles multi-step tasks", "Requests approvals"]),
        ("To Agentic Ops", ["Cross-system orchestration", "Human oversight", "Measured business outcomes"]),
    ],
    "Agents succeed only with guardrails and process ownership."
)

two_col(
    "OpenClaw Practical Pattern",
    "Architecture Elements",
    ["Tool-enabled agent loop", "State and memory handling", "Role-based action boundaries"],
    "Control Elements",
    ["Approval gates for external actions", "Audit trail for each run", "Fallback to human operator"],
    "Present as practical blueprint for controlled autonomy."
)

cards3(
    "Enterprise Pattern 1: Copilots",
    [
        ("Best Fit", ["Personal productivity", "Drafting and summarization", "Internal assistance"]),
        ("Strength", ["Fast deployment", "Low integration effort", "Quick visible wins"]),
        ("Constraint", ["Shallow process impact", "Harder ROI attribution", "Can create shadow AI"]),
    ],
    "Copilots are the on-ramp, not the destination."
)

cards3(
    "Enterprise Pattern 2: RAG",
    [
        ("Best Fit", ["Policy-heavy Q&A", "Knowledge retrieval", "Support use cases"]),
        ("Strength", ["Grounded responses", "Better traceability", "Faster ramp-up than full automation"]),
        ("Constraint", ["Dependent on data quality", "Index governance needed", "Permission model is critical"]),
    ],
    "RAG quality follows content governance quality."
)

cards3(
    "Enterprise Pattern 3: Agentic Workflows",
    [
        ("Best Fit", ["Repeatable multi-step operations", "Cross-system tasks", "High-volume back office flows"]),
        ("Strength", ["Largest value potential", "End-to-end cycle reduction", "Process-level learning"]),
        ("Constraint", ["Higher governance burden", "Needs observability", "Requires robust exception handling"]),
    ],
    "Start narrow and scale only after reliability gates are met."
)

two_col(
    "AI Operating Model",
    "Core Roles",
    ["Business owner", "AI platform lead", "Risk/compliance partner"],
    "Management Cadence",
    ["Monthly risk-review", "Quarterly value portfolio review", "Clear stop/scale criteria"],
    "Treat AI as a business capability, not a side project."
)

cards3(
    "EU + Swiss Governance",
    [
        ("EU AI Act", ["Phased obligations 2025-2027", "High-risk and transparency duties", "GPAI obligations active"]),
        ("Switzerland", ["FADP already applies", "Convention-aligned implementation", "Sector-specific guidance evolving"]),
        ("Action", ["Implement controls now", "Maintain legal watchlist", "Document oversight decisions"]),
    ],
    "Translate legal uncertainty into practical controls and accountability."
)

two_col(
    "Human Oversight Model",
    "Human-in-the-Loop",
    ["Legal/financial commitments", "Customer-impact decisions", "Sensitive communications"],
    "Human-on-the-Loop",
    ["Low-risk repetitive tasks", "Exception monitoring", "Threshold-based intervention"],
    "Define escalation triggers before go-live."
)

cards3(
    "AI Quality Management",
    [
        ("Pre-Prod", ["Golden datasets", "Task-specific benchmarks", "Failure mode tests"]),
        ("Post-Release", ["Regression suite", "Drift monitoring", "Model update playbooks"]),
        ("Metrics", ["Accuracy and latency", "Error budget", "Human correction rate"]),
    ],
    "Manage models with software-grade QA discipline."
)

cards3(
    "Security and Trust",
    [
        ("Threats", ["Prompt injection", "Data leakage", "Identity/deepfake fraud"]),
        ("Controls", ["Input/output filtering", "Least-privilege tools", "Runtime policy checks"]),
        ("Response", ["Incident playbooks", "Traceable logs", "Fast rollback capability"]),
    ],
    "Position trust as a technical and operating requirement."
)

two_col(
    "Data Readiness Checklist",
    "Must Have",
    ["Data classification", "Access and retention policy", "Metadata and lineage"],
    "Common Gaps",
    ["Stale documents", "Permission mismatches", "Unowned data domains"],
    "Data readiness is typically the #1 bottleneck for RAG programs."
)

cards3(
    "Build vs Buy: 6-Gate Decision",
    [
        ("Strategic", ["Differentiation", "IP sensitivity", "Reversibility"]),
        ("Operational", ["Integration depth", "Talent availability", "Time-to-value"]),
        ("Economic", ["24-month TCO", "Run-cost volatility", "Vendor dependence risk"]),
    ],
    "Recommendation: buy commodity, build differentiators, keep model portability."
)

timeline(
    "30/60/90 Day Adoption Roadmap",
    [
        ("Day 30", "Select top 3 use cases and governance baseline"),
        ("Day 60", "Run two secure pilots with KPI instrumentation"),
        ("Day 90", "Scale one winner and formalize operating model"),
    ],
    "Optional German: Nächster Schritt: Entscheidung in 2 Wochen."
)

cards3(
    "Use-Case Blueprint 1: Customer Service",
    [
        ("Scope", ["AI triage", "Drafted responses", "Knowledge retrieval"]),
        ("Controls", ["PII masking", "Escalation to human", "Confidence thresholds"]),
        ("KPIs", ["First response time", "Deflection rate", "CSAT movement"]),
    ],
    "Strong first use case due high volume and measurable outcomes."
)

cards3(
    "Use-Case Blueprint 2: Sales / RFP",
    [
        ("Scope", ["RAG-powered proposal drafting", "Evidence retrieval", "Versioned collaboration"]),
        ("Controls", ["Legal/compliance approval", "Source citation", "Final human sign-off"]),
        ("KPIs", ["Cycle-time reduction", "Bid quality", "Win-rate contribution"]),
    ],
    "Excellent fit for multilingual market contexts in Switzerland."
)

cards3(
    "Use-Case Blueprint 3: Procure-to-Pay",
    [
        ("Scope", ["Invoice extraction", "Policy checks", "Exception routing"]),
        ("Controls", ["Segregation of duties", "Threshold approvals", "Audit logs"]),
        ("KPIs", ["Processing time", "Error rate", "Audit exceptions"]),
    ],
    "A high-confidence path to measurable back-office ROI."
)

title_dark(
    "Executive Decisions Required This Month",
    "1) Approve 3 use cases  2) Approve governance baseline  3) Approve portable multi-model architecture",
    "Close with decision asks and ownership commitments."
)

out = "/root/.openclaw/workspace/AI_Executive_Briefing_Zurich_2026_v2.pptx"
prs.save(out)
print(out)
print("slides", len(prs.slides))
